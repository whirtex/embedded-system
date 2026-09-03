"""Deck do Seminario 1: trabalhos relacionados. Azul e branco."""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

SAIDA = "seminario-1-trabalhos-relacionados.pptx"

AZUL_FUNDO = RGBColor(0x0B, 0x33, 0x55)   # painel e capa
AZUL = RGBColor(0x1B, 0x6F, 0xB5)         # acento
AZUL_CLARO = RGBColor(0x7F, 0xB4, 0xDD)   # numerais e apoio sobre azul
AZUL_TENUE = RGBColor(0xE8, 0xF1, 0xF8)   # faixas sobre branco
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
TINTA = RGBColor(0x0F, 0x24, 0x36)        # texto sobre branco
CINZA = RGBColor(0x5B, 0x72, 0x85)

FONTE = "Arial"
L, A = Inches(13.333), Inches(7.5)
MARGEM = Inches(0.85)
UTIL = L - 2 * MARGEM

PAINEL = Inches(4.55)
DIR_X = PAINEL + Inches(0.75)
DIR_W = L - DIR_X - MARGEM


def bloco(slide, x, y, cx, cy, cor):
    forma = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cx, cy)
    forma.fill.solid()
    forma.fill.fore_color.rgb = cor
    forma.line.fill.background()
    forma.shadow.inherit = False
    forma.text_frame.text = ""
    return forma


def texto(slide, x, y, cx, cy, conteudo, tamanho, cor=TINTA, negrito=False,
          espaco=1.15, alinha=PP_ALIGN.LEFT, espacamento_letras=None):
    caixa = slide.shapes.add_textbox(x, y, cx, cy)
    tf = caixa.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    for i, linha in enumerate(conteudo if isinstance(conteudo, list) else [conteudo]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alinha
        p.line_spacing = espaco
        if i:
            p.space_before = Pt(9)
        r = p.add_run()
        r.text = linha
        r.font.size = Pt(tamanho)
        r.font.bold = negrito
        r.font.color.rgb = cor
        r.font.name = FONTE
        if espacamento_letras:
            r.font._rPr.set("spc", str(int(espacamento_letras * 100)))
    return caixa


def rotulo(slide, x, y, cx, conteudo, cor=AZUL):
    """Etiqueta pequena, maiúscula e espaçada."""
    return texto(slide, x, y, cx, Inches(0.28), conteudo.upper(), 12, cor,
                 negrito=True, espacamento_letras=1.4)


def topicos(slide, x, y, cx, itens, tamanho=16, cor=TINTA, marcador="1B6FB5"):
    caixa = slide.shapes.add_textbox(x, y, cx, Inches(0.4))
    tf = caixa.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    recuo = Inches(0.3)
    folga = 12 if len(itens) <= 3 else 8
    for i, item in enumerate(itens):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.22
        if i:
            p.space_before = Pt(folga)
        pPr = p._pPr
        pPr.set("marL", str(recuo))
        pPr.set("indent", str(-recuo))
        clr = pPr.makeelement(qn("a:buClr"), {})
        clr.append(pPr.makeelement(qn("a:srgbClr"), {"val": marcador}))
        pPr.append(clr)
        pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"}))
        pPr.append(pPr.makeelement(qn("a:buChar"), {"char": "–"}))
        r = p.add_run()
        r.text = item
        r.font.size = Pt(tamanho)
        r.font.color.rgb = cor
        r.font.name = FONTE
    return caixa


def altura_topicos(itens, largura_pol, tamanho=16):
    """Estimativa da altura ocupada, para posicionar o bloco seguinte."""
    por_linha = max(int((largura_pol - 0.3) * 148 / tamanho), 20)
    linhas = sum(max(1, -(-len(i) // por_linha)) for i in itens)
    folga = 0.17 if len(itens) <= 3 else 0.12
    return Inches(linhas * tamanho * 0.0185 + (len(itens) - 1) * folga)


def branco(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bloco(s, 0, 0, L, A, BRANCO)
    return s


def cabecalho(slide, etiqueta, titulo, linha=None):
    rotulo(slide, MARGEM, Inches(0.72), UTIL, etiqueta)
    texto(slide, MARGEM, Inches(1.15), UTIL, Inches(0.8), titulo, 38, TINTA, negrito=True)
    bloco(slide, MARGEM, Inches(1.95), Inches(1.9), Inches(0.055), AZUL)
    if linha:
        texto(slide, MARGEM, Inches(2.28), UTIL, Inches(0.4), linha, 19, CINZA)


def rodape(slide, direita, esquerda="Seminário 1 · Trabalhos relacionados"):
    y = A - Inches(0.6)
    bloco(slide, MARGEM, y - Inches(0.18), UTIL, Emu(9525), AZUL_TENUE)
    texto(slide, MARGEM, y, UTIL * 0.7, Inches(0.3), esquerda, 11, CINZA)
    texto(slide, MARGEM + UTIL * 0.7, y, UTIL * 0.3, Inches(0.3), direita, 11,
          CINZA, alinha=PP_ALIGN.RIGHT)


# ------------------------------------------------------------------ 1. capa
def capa(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bloco(s, 0, 0, L, A, AZUL_FUNDO)
    bloco(s, 0, 0, L, Inches(0.22), AZUL)

    rotulo(s, MARGEM, Inches(1.7), UTIL, "Seminário 1 · Trabalhos relacionados", AZUL_CLARO)
    texto(s, MARGEM, Inches(2.35), Inches(11.0), Inches(2.2),
          "Manutenção preditiva de aparelhos de ar-condicionado em operação contínua",
          44, BRANCO, negrito=True, espaco=1.08)
    bloco(s, MARGEM, Inches(4.9), Inches(1.9), Inches(0.055), AZUL_CLARO)
    texto(s, MARGEM, Inches(5.3), UTIL, Inches(0.9),
          ["Igor Costa · Jorge Alves · Ian Dias · Davi Ito",
           "Ibmec, Rio de Janeiro · IBM3118 Sistemas Embarcados · 2026.2 · Grupo 3"],
          18, AZUL_CLARO, espaco=1.35)


# ------------------------------------------------------------- 2. motivação
def motivacao(prs):
    s = branco(prs)
    cabecalho(s, "Motivação", "Começou pelas áreas com pacientes",
              "Aparelhos instalados nos ambientes onde há paciente o tempo todo.")

    itens = [
        ("Enfermarias e quartos",
         "Ligado dia e noite. Parar significa remanejar quem está internado."),
        ("UTIs e leitos críticos",
         "O paciente não pode ser removido, e a temperatura faz parte do cuidado."),
        ("Pronto-socorro",
         "Ocupação contínua. Quase nunca é desligado para manutenção."),
    ]

    painel_y, painel_h = Inches(2.9), Inches(2.2)
    bloco(s, MARGEM, painel_y, UTIL, painel_h, AZUL_TENUE)

    # margem interna igual nos quatro lados do retângulo
    recuo, vao = Inches(0.45), Inches(0.5)
    larg = (UTIL - recuo * 2 - vao * 2) / 3
    conteudo_y = painel_y + recuo
    for i, (titulo, corpo) in enumerate(itens):
        x = MARGEM + recuo + i * (larg + vao)
        if i:
            bloco(s, x - vao / 2, conteudo_y, Emu(12700), painel_h - recuo * 2, AZUL_CLARO)
        texto(s, x, conteudo_y, larg, Inches(0.4), titulo, 19, TINTA, negrito=True)
        texto(s, x, conteudo_y + Inches(0.52), larg, Inches(0.9),
              corpo, 15, CINZA, espaco=1.3)

    bloco(s, MARGEM, Inches(5.5), UTIL, Inches(1.08), AZUL_FUNDO)
    texto(s, MARGEM + Inches(0.45), Inches(5.77), UTIL - Inches(0.9), Inches(0.6),
          "O desgaste aparece antes da parada, na vibração do compressor, "
          "na corrente do motor e na temperatura de operação.",
          17, BRANCO, **{"espaco": 1.3})
    rodape(s, "2 / 11")


# ---------------------------------------------------------- 3. outros locais
def outros_locais(prs):
    s = branco(prs)
    cabecalho(s, "Alcance", "O mesmo problema aparece fora do hospital",
              "O critério é operação ininterrupta somada a uma parada que custa mais que o reparo.")

    locais = [
        ("Data centers e salas de servidores",
         "A refrigeração é o ponto único de falha. Sem ela, o hardware desliga por temperatura."),
        ("Laboratórios refrigerados",
         "Reagentes e amostras fora de faixa invalidam resultados e perdem-se em silêncio."),
        ("Centros de operações e emergência",
         "Salas do 190 e do 193 nunca desligam. Se o equipamento cai, o atendimento cai junto."),
        ("Aeroportos e rodoviárias",
         "Terminais abertos o tempo todo. A manutenção precisa acontecer sem interromper a operação."),
    ]
    larg = (UTIL - Inches(0.9)) / 2
    for i, (titulo, corpo) in enumerate(locais):
        x = MARGEM + (i % 2) * (larg + Inches(0.9))
        y = Inches(3.05) + (i // 2) * Inches(1.75)
        bloco(s, x, y + Inches(0.06), Inches(0.055), Inches(1.2), AZUL)
        texto(s, x + Inches(0.35), y, larg - Inches(0.5), Inches(0.45), titulo, 20, TINTA, negrito=True)
        texto(s, x + Inches(0.35), y + Inches(0.55), larg - Inches(0.5), Inches(0.9),
              corpo, 16, CINZA, espaco=1.25)

    texto(s, MARGEM, Inches(6.4), UTIL, Inches(0.4),
          "Muda a consequência da parada. O equipamento e o método continuam os mesmos.",
          17, CINZA)
    rodape(s, "3 / 11")


# -------------------------------------------------- 4 a 9. uma por referência
REFERENCIAS = [
    dict(
        camada="Panorama",
        autor="Meitz et al.", ano="2025",
        problema="Revisão estruturada e desafios em aberto da manutenção preditiva na Indústria 4.0",
        estudou=[
            "Revisão de 249 publicações, organizadas em nove categorias",
            "Monitoramento de condição: o que medir no equipamento",
            "Tratamento de dados: limpar, rotular e organizar",
            "Detecção de falhas, degradação e prognóstico",
            "Avaliação e planejamento da manutenção",
        ],
        levamos=[
            "As etapas dele viram a estrutura do nosso projeto",
            "Ficamos em monitoramento, dados, detecção e avaliação",
            "Prognóstico e vida útil restante ficam declaradamente fora",
            "Avaliação é etapa própria, não apêndice do modelo",
        ],
        fonte=["Computers & Industrial Engineering", "v. 206, art. 111193, 2025",
               "DOI 10.1016/j.cie.2025.111193"],
    ),
    dict(
        camada="Monitoramento de condição",
        autor="Yousuf et al.", ano="2024",
        problema="Monitoramento de condição e detecção de falhas em motor de indução CA",
        estudou=[
            "Motor de indução CA instrumentado com temperatura, vibração, corrente, tensão e velocidade",
            "Aquisição em Arduino, com alarme local e notificação por GSM",
            "Proteção automática por relé e histórico na plataforma IoT Blynk",
        ],
        levamos=[
            "Fechar a cadeia inteira, do sensor até a ação",
            "Temperatura, vibração e corrente como lista de partida",
            "Relé como resposta automática à falha",
        ],
        fonte=["Measurement and Control", "v. 57, n. 8, 2024",
               "DOI 10.1177/00202940241231473"],
    ),
    dict(
        camada="Monitoramento de condição",
        autor="Mohammed et al.", ano="2023",
        problema="Manutenção preditiva de motores elétricos com IoT industrial e aprendizado de máquina",
        estudou=[
            "Raspberry Pi coletando vibração, corrente e temperatura",
            "Transmissão por MQTT para servidor em nuvem",
            "Cinco algoritmos supervisionados sobre falhas induzidas, com Random Forest à frente",
        ],
        levamos=[
            "MQTT separa aquisição de análise e libera o ESP32",
            "Falha induzida gera dados sem histórico prévio",
            "Coleta a cada segundo não preserva a vibração",
        ],
        fonte=["J. Européen des Systèmes Automatisés", "v. 56, n. 4, 2023",
               "DOI 10.18280/jesa.560414"],
    ),
    dict(
        camada="Detecção de falhas",
        autor="Kolok et al.", ano="2025",
        problema="Manutenção preditiva de baixo custo baseada em vibração",
        estudou=[
            "ESP32 com sensores MEMS de vibração e acústico",
            "RMS no domínio do tempo e FFT no da frequência",
            "Isolation Forest treinado apenas com operação saudável",
        ],
        levamos=[
            "Dispensa exemplos de falha, nossa maior limitação",
            "Extração de característica cabe na borda",
            "Calibração por equipamento é obrigatória",
        ],
        fonte=["Sensors", "v. 25, art. 6610, 2025", "DOI 10.3390/s25216610"],
    ),
    dict(
        camada="Tratamento de dados e avaliação",
        autor="Gupta et al.", ano="2023",
        problema="Manutenção preditiva de esteiras de bagagem de aeroporto com IoT",
        estudou=[
            "Vibração por IoT em oito esteiras idênticas em operação real",
            "Sem histórico até a falha, com limpeza apoiada em RMS",
            "Rótulos extraídos de registros de manutenção em texto",
            "Quatro classificadores comparados, com Random Forest à frente",
        ],
        levamos=[
            "É o caso mais próximo do nosso",
            "Limpar ruído e rotular consomem a maior parte do trabalho",
            "Registro escrito de manutenção pode virar rótulo",
        ],
        fonte=["Computers & Industrial Engineering", "v. 177, art. 109033, 2023",
               "DOI 10.1016/j.cie.2023.109033"],
    ),
    dict(
        camada="Interpretação do alerta",
        autor="Tormos et al.", ano="2026",
        problema="Detecção de anomalia e explicação de alertas em ar-condicionado de frota de ônibus",
        estudou=[
            "Ar-condicionado de ônibus urbanos monitorado por sensores e telemetria",
            "Sensores mapeados sobre o ciclo de refrigeração antes de modelar",
            "Isolation Forest aprende a operação normal, sem rótulo de falha",
            "Kernel SHAP aponta qual variável causou cada desvio",
        ],
        levamos=[
            "Aprender o normal dispensa exemplos de falha",
            "O alerta indica qual medição disparou o alarme",
            "A explicação roda no servidor, não no ESP32",
        ],
        fonte=["Algorithms", "v. 19, n. 7, art. 586, 2026",
               "DOI 10.3390/a19070586", "Acesso aberto, CC BY 4.0"],
    ),
]


# ------------------------------------------------- 5. as etapas do Meitz
def etapas(prs):
    s = branco(prs)
    cabecalho(s, "A estrutura do projeto", "As etapas que o Meitz define",
              "Quatro delas organizam o nosso trabalho. A quinta é acréscimo nosso.")

    dentro = [
        ("Monitoramento\nde condição", ["[2] Yousuf", "[3] Mohammed"]),
        ("Tratamento\nde dados", ["[5] Gupta"]),
        ("Detecção\nde falhas", ["[4] Kolok"]),
        ("Avaliação", ["[1] Meitz", "[5] Gupta"]),
        ("Interpretação\ndo alerta", ["[6] Tormos"]),
    ]
    vao = Inches(0.24)
    larg = (UTIL - vao * 4) / 5
    topo, alt = Inches(2.95), Inches(1.2)
    for i, (nome, refs) in enumerate(dentro):
        x = MARGEM + i * (larg + vao)
        bloco(s, x, topo, larg, alt, AZUL_FUNDO)
        texto(s, x + Inches(0.2), topo + Inches(0.26), larg - Inches(0.4), Inches(0.8),
              nome.split("\n"), 15, BRANCO, negrito=True, espaco=1.15)
        texto(s, x, topo + alt + Inches(0.16), larg, Inches(0.5),
              refs, 11, AZUL, negrito=True, espaco=1.25)
        if i < 4:
            texto(s, x + larg, topo + Inches(0.38), vao, Inches(0.4), "\u203a", 22,
                  AZUL, negrito=True, alinha=PP_ALIGN.CENTER)

    bloco(s, MARGEM, Inches(5.05), UTIL, Emu(9525), AZUL_TENUE)
    rotulo(s, MARGEM, Inches(5.32), UTIL, "Fora do escopo do protótipo", CINZA)

    fora = ["Modelagem de degradação", "Prognóstico", "Planejamento da manutenção"]
    larg2 = (UTIL - vao * 2) / 3
    for i, nome in enumerate(fora):
        x = MARGEM + i * (larg2 + vao)
        bloco(s, x, Inches(5.72), larg2, Inches(0.6), AZUL_TENUE)
        texto(s, x + Inches(0.25), Inches(5.88), larg2 - Inches(0.5), Inches(0.4),
              nome, 15, CINZA, negrito=True)

    texto(s, MARGEM, Inches(6.52), UTIL, Inches(0.4),
          "As três exigem histórico de falhas validadas, que o protótipo não terá.",
          15, CINZA)
    rodape(s, "5 / 11")


def slide_referencia(prs, ref, indice, pagina):
    s = branco(prs)
    bloco(s, 0, 0, PAINEL, A, AZUL_FUNDO)

    largura_painel = PAINEL - MARGEM - Inches(0.35)
    corpo_autor = 29 if len(ref["autor"]) <= 14 else 24 if len(ref["autor"]) <= 20 else 21
    texto(s, MARGEM, Inches(0.8), Inches(2.0), Inches(1.2), f"{indice:02d}", 60, AZUL, negrito=True)
    rotulo(s, MARGEM, Inches(2.0), Inches(3.2), ref["camada"], AZUL_CLARO)
    texto(s, MARGEM, Inches(2.5), largura_painel, Inches(1.6),
          ref["autor"], corpo_autor, BRANCO, negrito=True, espaco=1.15)

    # o ano acompanha a altura ocupada pelo nome
    por_linha = int(largura_painel / 914400 * 148 / corpo_autor)
    linhas_autor = max(1, -(-len(ref["autor"]) // por_linha))
    y_ano = Inches(2.5) + Inches(linhas_autor * corpo_autor * 0.019 + 0.2)
    texto(s, MARGEM, y_ano, Inches(2.0), Inches(0.5), ref["ano"], 22, AZUL_CLARO, negrito=True)
    bloco(s, MARGEM, Inches(4.3), Inches(1.4), Inches(0.045), AZUL)
    texto(s, MARGEM, Inches(4.7), largura_painel, Inches(1.4),
          ref["fonte"], 12, AZUL_CLARO, espaco=1.35)

    texto(s, DIR_X, Inches(0.95), DIR_W, Inches(0.9), ref["problema"], 20, TINTA,
          negrito=True, espaco=1.2)

    y = Inches(2.1)
    rotulo(s, DIR_X, y, DIR_W, "O que o trabalho estudou")
    topicos(s, DIR_X, y + Inches(0.45), DIR_W, ref["estudou"])

    y2 = y + Inches(0.45) + altura_topicos(ref["estudou"], DIR_W / 914400) + Inches(0.32)
    bloco(s, DIR_X, y2, DIR_W, Emu(22860), AZUL)
    rotulo(s, DIR_X, y2 + Inches(0.3), DIR_W, "O que levamos para o projeto")
    topicos(s, DIR_X, y2 + Inches(0.75), DIR_W, ref["levamos"])

    texto(s, MARGEM, A - Inches(0.62), largura_painel, Inches(0.3),
          f"{pagina} / 11", 11, AZUL_CLARO)


# ----------------------------------------------------------- 10. agradecimento
def agradecimento(prs):
    """Mesma composição do slide final do deck de colmeias, na paleta azul."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bloco(s, 0, 0, L, A, AZUL_FUNDO)

    centro = dict(alinha=PP_ALIGN.CENTER)
    texto(s, Inches(0.60), Inches(2.72), Inches(12.13), Inches(1.10),
          "Obrigado!", 60, BRANCO, negrito=True, **centro)
    texto(s, Inches(0.60), Inches(3.86), Inches(12.13), Inches(0.50),
          "Perguntas e Discussão", 20, AZUL_CLARO, **centro)
    texto(s, Inches(0.60), Inches(4.95), Inches(12.13), Inches(0.30),
          "GRUPO 3  ·  IBMEC RIO DE JANEIRO", 11, AZUL_CLARO, negrito=True, **centro)
    texto(s, Inches(0.60), Inches(5.30), Inches(12.13), Inches(0.40),
          "Igor Costa   ·   Jorge Alves   ·   Ian Dias   ·   Davi Ito",
          15, BRANCO, negrito=True, **centro)
    texto(s, Inches(0.60), Inches(7.02), Inches(12.13), Inches(0.28),
          "Manutenção preditiva de ar-condicionado  ·  IBM3118 Sistemas Embarcados (2026.2)",
          9, AZUL_CLARO, **centro)


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = L, A

    capa(prs)
    motivacao(prs)
    outros_locais(prs)
    paginas = [4, 6, 7, 8, 9, 10]
    for i, (ref, pagina) in enumerate(zip(REFERENCIAS, paginas), start=1):
        slide_referencia(prs, ref, i, pagina)
        if i == 1:
            etapas(prs)
    agradecimento(prs)

    prs.save(SAIDA)
    print(f"{SAIDA}: {len(prs.slides._sldIdLst)} slides")


if __name__ == "__main__":
    main()
